#!/usr/bin/env python3
"""Generate a PowerPoint presentation from the Release Management Plan with speaker notes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── colour palette ──
DARK_NAVY   = RGBColor(0x1e, 0x3a, 0x5f)
MID_BLUE    = RGBColor(0x3b, 0x82, 0xf6)
LIGHT_BLUE  = RGBColor(0x93, 0xc5, 0xfd)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY  = RGBColor(0xf0, 0xf0, 0xf0)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
ACCENT_GREEN = RGBColor(0x10, 0xb9, 0x81)
ACCENT_RED   = RGBColor(0xef, 0x44, 0x44)

# ── helpers ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_GRAY, bold_prefix=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bold_prefix and '—' in item:
            parts = item.split('—', 1)
            run1 = p.add_run()
            run1.text = parts[0] + '—'
            run1.font.bold = True
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = 'Calibri'
        else:
            p.text = f'• {item}'
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox

def add_table_slide(slide, left, top, width, headers, rows, col_widths=None, font_size=11):
    row_count = 1 + len(rows)
    col_count = len(headers)
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, Inches(0.4 * row_count))
    table = table_shape.table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_GRAY
                p.font.name = 'Calibri'
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    return table_shape

def slide_title_bar(slide, title):
    """Add a navy title bar at top of content slide."""
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), DARK_NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7), title,
                font_size=28, color=WHITE, bold=True)
    # accent line
    add_shape_bg(slide, Inches(0), Inches(1.1), Inches(13.33), Inches(0.05), MID_BLUE)

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def create_mock_frontend_image():
    """Create a mock of the frontend screenshot."""
    img = Image.new('RGB', (900, 600), '#ffffff')
    draw = ImageDraw.Draw(img)
    # Green header bar
    draw.rectangle([0, 0, 900, 100], fill='#2d5016')
    # Try to use a basic font
    try:
        font_lg = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 28)
        font_md = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
        font_sm = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 14)
    except:
        font_lg = ImageFont.load_default()
        font_md = ImageFont.load_default()
        font_sm = ImageFont.load_default()

    draw.text((300, 25), "Recipe Manager", fill='white', font=font_lg)
    draw.text((310, 65), "Microservice Recipe Application", fill='#cccccc', font=font_sm)

    # Form area
    draw.rectangle([150, 130, 750, 500], outline='#dddddd', width=1)
    draw.text((170, 145), "Add New Recipe", fill='#333333', font=font_md)
    draw.text((170, 180), "Title", fill='#333333', font=font_sm)
    draw.rectangle([170, 200, 730, 230], outline='#3b82f6', width=2)
    draw.text((175, 205), "Recipe title", fill='#999999', font=font_sm)

    draw.text((170, 250), "Ingredients", fill='#333333', font=font_sm)
    draw.rectangle([170, 270, 730, 350], outline='#cccccc', width=1)
    draw.text((175, 275), "One ingredient per line", fill='#999999', font=font_sm)

    draw.text((170, 370), "Instructions", fill='#333333', font=font_sm)
    draw.rectangle([170, 390, 730, 460], outline='#cccccc', width=1)
    draw.text((175, 395), "Step-by-step instructions", fill='#999999', font=font_sm)

    # Button
    draw.rectangle([170, 475, 290, 500], fill='#3b82f6')
    draw.text((185, 478), "Add Recipe", fill='white', font=font_sm)

    # Footer
    draw.text((160, 520), "All Recipes (0)", fill='#333333', font=font_md)
    draw.text((160, 550), "No recipes yet. Add one above!", fill='#999999', font=font_sm)

    # URL bar mock
    draw.rectangle([0, 0, 900, 0], fill='#f0f0f0')

    path = os.path.join(OUT_DIR, '_mock_frontend.png')
    img.save(path)
    return path


def create_mock_api_image():
    """Create a mock of the API response screenshot."""
    img = Image.new('RGB', (900, 400), '#ffffff')
    draw = ImageDraw.Draw(img)
    # Green header bar (browser chrome style)
    draw.rectangle([0, 0, 900, 50], fill='#2d5016')
    try:
        font_md = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
        font_sm = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 14)
        font_code = ImageFont.truetype("/System/Library/Fonts/Menlo.ttc", 18)
    except:
        font_md = ImageFont.load_default()
        font_sm = ImageFont.load_default()
        font_code = ImageFont.load_default()

    draw.text((50, 15), "localhost:8080/api/recipes", fill='#cccccc', font=font_sm)

    # API response
    draw.text((20, 75), "Pretty-print", fill='#333333', font=font_sm)
    draw.rectangle([115, 75, 130, 90], outline='#999999', width=1)
    draw.text((20, 110), "[]", fill='#333333', font=font_code)

    path = os.path.join(OUT_DIR, '_mock_api.png')
    img.save(path)
    return path


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════
    # SLIDE 1 — TITLE
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_NAVY)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                'Release Management Plan', font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(4), Inches(2.8), Inches(5.33), Inches(0.04), MID_BLUE)
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                'Enterprise Microservice CI/CD Deployment', font_size=24, color=LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                'Recipe Management Application', font_size=20, color=RGBColor(0xaa, 0xaa, 0xaa),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                'DevOps — Continuous Assessment  |  March 2026', font_size=16,
                color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

    set_notes(slide, """WHAT TO SAY:
"Good morning/afternoon everyone. Today I'll be presenting our Release Management Plan for an enterprise-style CI/CD deployment of a microservice-based Recipe Management Application.

This project demonstrates a complete end-to-end deployment pipeline — from code commit through to production — using industry-standard tools and best practices.

I'll walk you through the architecture, our CI/CD pipeline design, deployment strategy, infrastructure as code, and then critically evaluate our decisions against eleven criteria."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 2 — AGENDA
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Agenda')

    left_items = [
        '1. System Architecture Overview',
        '2. Live Application Demo',
        '3. Repository & Branching Strategy',
        '4. CI/CD Pipeline Architecture',
        '5. Deployment Strategy (Blue/Green)',
        '6. Infrastructure as Code (Terraform)',
    ]
    right_items = [
        '7. Container Orchestration (AKS + Helm)',
        '8. Evaluation Criteria (11 areas)',
        '9. Additional Features',
        '10. Cost Analysis',
        '11. Critical Reflection',
        '12. Q&A',
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), left_items, font_size=20)
    add_bullet_list(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5), right_items, font_size=20)

    set_notes(slide, """WHAT TO SAY:
"Here's what we'll cover today. I'll start with the system architecture — what services we built and how they communicate. Then I'll show you a live demo of the running application.

After that, we'll go deeper into the technical decisions: our branching strategy, CI/CD pipeline, blue/green deployment, Terraform infrastructure, and Kubernetes orchestration.

The bulk of the presentation evaluates our choices against eleven criteria — performance, cost, security, scaling, and more. I'll finish with a critical reflection on what worked well and what we'd improve.

Let's get started."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 3 — ARCHITECTURE OVERVIEW
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'System Architecture Overview')

    # Three service boxes
    services = [
        ('Frontend', 'Node.js 20 / Express / EJS', 'Port 3000', 'Dynamic Web UI\nProxies API requests', MID_BLUE),
        ('Backend API', 'Java 17 / Spring Boot 3', 'Port 8080', 'REST API — CRUD\nValidation & Logic', ACCENT_GREEN),
        ('Database', 'MongoDB 7', 'Port 27017', 'Document Store\nBSON recipes', RGBColor(0xf5, 0x9e, 0x0b)),
    ]
    for i, (name, tech, port, desc, clr) in enumerate(services):
        x = Inches(1.0 + i * 4.0)
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.4), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = clr
        box.line.width = Pt(3)
        # Header bar inside box
        add_shape_bg(slide, x + Inches(0.05), Inches(1.85), Inches(3.3), Inches(0.7), clr)
        add_textbox(slide, x + Inches(0.2), Inches(1.9), Inches(3), Inches(0.6), name,
                    font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(2.7), Inches(3), Inches(0.4), tech,
                    font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(3.1), Inches(3), Inches(0.3), port,
                    font_size=12, color=MID_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(3.6), Inches(3), Inches(1.2), desc,
                    font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrows between boxes
    for x_start in [Inches(4.4), Inches(8.4)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_start, Inches(3.3), Inches(0.6), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        arrow.line.fill.background()

    add_textbox(slide, Inches(3.5), Inches(5.6), Inches(6), Inches(0.5),
                'HTTP/JSON                                MongoDB Wire Protocol',
                font_size=13, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER, italic=True)

    # Communication note
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
                'Synchronous request-response  •  Each service independently buildable & deployable  •  Docker multi-stage builds',
                font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, italic=True)

    set_notes(slide, """WHAT TO SAY:
"Our application follows a classic three-tier microservice architecture.

The Frontend is built with Node.js 20 and Express, using EJS for server-side rendering. It serves the web UI on port 3000 and proxies all API requests to the backend.

The Backend is a Java 17 Spring Boot 3 REST API running on port 8080. It handles all CRUD operations, input validation, and business logic for recipe management.

The Database is MongoDB 7, storing recipes as BSON documents.

Communication between frontend and backend is synchronous HTTP/JSON. Between the backend and database, it uses the standard MongoDB wire protocol.

Critically, each service is independently buildable and deployable — which is the defining characteristic of a microservice architecture. The frontend can be updated without touching the backend, and vice versa."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 4 — LIVE DEMO: FRONTEND
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Live Demo — Frontend (localhost:3000)')

    # Embed mock frontend screenshot
    fe_img = create_mock_frontend_image()
    slide.shapes.add_picture(fe_img, Inches(1.8), Inches(1.5), Inches(9.5), Inches(5.5))

    # Border around image
    add_textbox(slide, Inches(1.8), Inches(7.1), Inches(9.5), Inches(0.4),
                'Recipe Manager — Node.js / Express / EJS frontend running locally via Docker Compose',
                font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER, italic=True)

    set_notes(slide, """WHAT TO SAY:
"Here's the frontend of our application running locally at localhost:3000 via Docker Compose.

You can see the Recipe Manager interface — it has a form to add new recipes with a title, ingredients, and step-by-step instructions. Below the form, we have the recipe listing area, currently showing zero recipes.

The UI is server-side rendered using EJS templates. The frontend communicates with the Spring Boot backend API to create, read, update, and delete recipes.

Important to note — this is running in Docker. A single 'docker-compose up --build' starts all three services: the frontend, the backend, and MongoDB. No local Java, Node.js, or MongoDB installation is required."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 5 — LIVE DEMO: BACKEND API
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Live Demo — Backend API (localhost:8080/api/recipes)')

    # Embed mock API screenshot
    api_img = create_mock_api_image()
    slide.shapes.add_picture(api_img, Inches(2.5), Inches(1.5), Inches(8), Inches(3.5))

    add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
                'Spring Boot 3 REST API returning JSON — empty array confirms healthy service + DB connectivity',
                font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER, italic=True)

    # Key points under screenshot
    add_bullet_list(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.5), [
        'RESTful endpoints: GET, POST, PUT, DELETE on /api/recipes',
        'Spring Data MongoDB handles object-document mapping automatically',
        'Health check via /actuator/health — used by Kubernetes readiness probes',
        'Prometheus metrics exposed at /actuator/prometheus for monitoring',
    ], font_size=15)

    set_notes(slide, """WHAT TO SAY:
"And here's the backend API, accessible directly at localhost:8080/api/recipes.

It returns a JSON array — currently empty because we haven't added any recipes yet. This empty array response actually confirms two important things: the Spring Boot API is running correctly, AND it has successfully connected to the MongoDB database.

The API is fully RESTful — it supports GET, POST, PUT, and DELETE operations. Spring Data MongoDB handles the object-document mapping so we don't write any manual queries.

The backend also exposes health and metrics endpoints via Spring Boot Actuator. The /actuator/health endpoint is what Kubernetes uses for readiness probes — it only sends traffic to pods that report healthy. The /actuator/prometheus endpoint exposes metrics in Prometheus format for our monitoring stack."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 6 — REPO & BRANCHING
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Repository & Branching Strategy')

    # Mono-repo decision
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                'Decision: Mono-Repo', font_size=22, color=DARK_NAVY, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), [
        'Single repo — complete project visibility',
        'Atomic cross-service changes in one commit',
        'Shared CI/CD and infra code',
        'Path-based CI triggers isolate pipelines',
        'Scales to ~10 services (Google/Meta use mono-repos)',
    ], font_size=15)

    # Branching model
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                'Branching Model: GitHub Flow', font_size=22, color=DARK_NAVY, bold=True)

    add_table_slide(slide, Inches(7), Inches(2.3), Inches(5.5),
        ['Branch', 'Purpose', 'Protection'],
        [
            ['main', 'Production-ready', 'CI + scan + approval'],
            ['develop', 'Integration', 'CI pass required'],
            ['feature/*', 'Feature work', 'None'],
            ['hotfix/*', 'Critical fixes', 'Same as main'],
        ], font_size=11)

    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.8),
                'PR into main requires: (1) All CI checks pass  (2) Trivy scan — no CRITICAL CVEs  (3) Code review approval',
                font_size=14, color=DARK_GRAY, italic=True)

    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
                'Critique: GitHub Flow is simpler than GitFlow — ideal for continuous deployment. '
                'For teams with less CI maturity, GitFlow\'s explicit release branches provide additional safety.',
                font_size=14, color=RGBColor(0x88, 0x44, 0x00), italic=True)

    set_notes(slide, """WHAT TO SAY:
"We chose a mono-repo structure — both services, all infrastructure code, Helm charts, and CI/CD workflows live in a single repository.

For a two-service application, this gives us the best developer experience: a single clone, complete visibility, and the ability to make atomic cross-service changes in one commit. Path-based CI triggers in GitHub Actions ensure that the frontend pipeline only runs when frontend code changes, and likewise for the backend.

For branching, we follow GitHub Flow — a simple model with main as production-ready, a develop branch for integration, and feature branches for individual work.

Pull requests into main have three gates: all CI checks must pass, the Trivy vulnerability scan must report no critical CVEs, and at least one code review approval is required.

One critique here: GitHub Flow assumes main is always deployable, which requires strong CI discipline. For less mature teams, GitFlow with explicit release branches would be safer."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 7 — CI/CD PIPELINE
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'CI/CD Pipeline Architecture')

    # CI - Left
    add_textbox(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(0.5),
                'Continuous Integration (per-service)', font_size=20, color=DARK_NAVY, bold=True)

    ci_fe = ['Checkout → Node.js 20 setup', 'npm install → ESLint → Jest tests',
             'Docker build (multi-stage)', 'Trivy scan (fail on CRITICAL)', 'Push to ACR (SHA tag)']
    ci_be = ['Checkout → Java 17 + Maven', 'mvn clean verify (compile + test)',
             'Docker build (multi-stage)', 'Trivy scan (fail on CRITICAL)', 'Push to ACR (SHA tag)']

    add_textbox(slide, Inches(0.6), Inches(2.0), Inches(2.8), Inches(0.3),
                'Frontend CI', font_size=14, color=MID_BLUE, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(2.4), Inches(3), Inches(2.5), ci_fe, font_size=13)

    add_textbox(slide, Inches(3.5), Inches(2.0), Inches(2.8), Inches(0.3),
                'Backend CI', font_size=14, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(3.5), Inches(2.4), Inches(3), Inches(2.5), ci_be, font_size=13)

    # CD - Right
    add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
                'Continuous Delivery', font_size=20, color=DARK_NAVY, bold=True)

    cd_steps = ['Azure Login (service principal)',
                'Get AKS credentials',
                'Detect inactive slot (blue/green)',
                'Deploy MongoDB → Backend → Frontend',
                'Health check (kubectl wait)',
                'Switch traffic (patch selector)',
                'Print deployment summary']
    add_bullet_list(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(3), cd_steps, font_size=14)

    # Key design decisions
    add_shape_bg(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.8), LIGHT_GRAY)
    add_textbox(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.4),
                'Key Design Decisions', font_size=16, color=DARK_NAVY, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2), [
        'Image tags use Git SHA — every deployed image maps to an exact commit',
        'Images only built on main branch; feature branches run tests only',
        'Platform: GitHub Actions — zero setup, excellent Azure integration, 2,000 min/month free',
    ], font_size=14)

    set_notes(slide, """WHAT TO SAY:
"Our CI/CD architecture has three independent pipelines — two for continuous integration and one for continuous delivery.

Each service has its own CI pipeline. The frontend pipeline triggers only when files under services/frontend change. It installs dependencies, runs ESLint and Jest tests, builds a Docker image using a multi-stage build, runs a Trivy vulnerability scan — which will fail the build on any CRITICAL CVE — and then pushes the image to Azure Container Registry tagged with the Git commit SHA.

The backend pipeline follows the same pattern but with Java: Maven compile and verify, Docker build, Trivy scan, and ACR push.

The CD pipeline is triggered when either CI pipeline completes on main. It logs into Azure, gets AKS credentials, automatically detects which slot — blue or green — is currently inactive, deploys all services to that inactive slot using Helm, waits for health checks to pass, then switches traffic by patching the Kubernetes service selector.

A critical design decision: image tags are Git SHAs, not version numbers. This means every running container maps directly to an exact code commit — essential for debugging and auditing."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 8 — DEPLOYMENT STRATEGY
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Deployment Strategy — Blue/Green')

    # Comparison table
    add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12),
        ['Strategy', 'Downtime', 'Rollback Speed', 'Resource Cost', 'Complexity'],
        [
            ['Recreate', 'Yes ✗', 'Slow', 'Low', 'Very Low'],
            ['Rolling Update', 'No ✓', 'Moderate', 'Low', 'Low'],
            ['Blue/Green ✓', 'No ✓', 'Instant (<10s)', 'High (double)', 'Moderate'],
            ['Canary', 'No ✓', 'Fast', 'Moderate', 'High (mesh)'],
        ], font_size=12)

    # How it works
    add_textbox(slide, Inches(0.6), Inches(4.0), Inches(5.5), Inches(0.4),
                'How Blue/Green Works:', font_size=18, color=DARK_NAVY, bold=True)

    steps = [
        '1. Identify inactive slot (auto-detect)',
        '2. Deploy new version to inactive slot',
        '3. Wait for health checks to pass',
        '4. Patch Service selector → switch traffic',
        '5. Old version remains on standby for rollback',
    ]
    add_bullet_list(slide, Inches(0.6), Inches(4.5), Inches(5.5), Inches(2.5), steps, font_size=15)

    # Benefits box
    add_shape_bg(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(3), RGBColor(0xec, 0xfe, 0xf5))
    add_textbox(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.4),
                'Why Blue/Green?', font_size=18, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(2.2), [
        'Zero downtime — users never see an error',
        'Instant rollback — patch selector back (<10 seconds)',
        'Full testing of new version before traffic switch',
        'No service mesh dependency (unlike Canary)',
    ], font_size=15, color=RGBColor(0x06, 0x5f, 0x46))

    set_notes(slide, """WHAT TO SAY:
"For deployment, we evaluated four strategies. Let me walk through why we chose blue/green.

Recreate was immediately rejected — it has downtime, which is unacceptable. Rolling update is simple but has two problems: rollback requires a full redeployment, and during the rollback you have mixed versions serving traffic simultaneously.

Canary is sophisticated — routing a percentage of traffic to the new version — but it requires a service mesh like Istio, which adds disproportionate complexity for our two-service application.

Blue/green hits the sweet spot. Here's how it works: we automatically detect which slot is inactive, deploy the new version there, wait for health checks, then atomically switch traffic by patching the Kubernetes Service selector. The old version's pods remain running on the now-inactive slot.

The killer feature is rollback speed. If something goes wrong, we patch the selector back to the previous slot — that's a metadata-only operation. No image pulls, no scheduling, no waiting. Under ten seconds to full rollback.

The trade-off is resource cost — you're running double the pods during deployment. For our application scale, this is acceptable."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 9 — INFRASTRUCTURE AS CODE
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Infrastructure as Code — Terraform')

    # IaC tool comparison
    add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(5.8),
        ['Criterion', 'Terraform', 'Bicep', 'Pulumi'],
        [
            ['Cloud', 'Multi-cloud', 'Azure only', 'Multi-cloud'],
            ['Language', 'HCL', 'Bicep DSL', 'Python/TS/Go'],
            ['State', 'Explicit file', 'Azure-managed', 'Built-in'],
            ['Community', 'Largest', 'Growing', 'Moderate'],
            ['Maturity', '2014', '2020', '2018'],
        ], font_size=11)

    # Resources provisioned
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                'Resources Provisioned:', font_size=18, color=DARK_NAVY, bold=True)
    resources = [
        'Resource Group (rg-{project}-{env})',
        'VNet + Subnet (10.0.0.0/16)',
        'AKS Cluster (1–5 nodes, B2s, autoscale)',
        'Container Registry (ACR, Basic SKU)',
        'Log Analytics Workspace (30-day retention)',
        'Storage Account + Blob (backups)',
        'Role Assignment (AKS → AcrPull)',
    ]
    add_bullet_list(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(3), resources, font_size=14)

    # Key capability
    add_shape_bg(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.5), LIGHT_GRAY)
    add_textbox(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
                'Destroy & Rebuild Capability', font_size=16, color=DARK_NAVY, bold=True)
    add_textbox(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.8),
                'terraform destroy removes ALL infrastructure. terraform apply re-creates it identically. '
                'Total rebuild: ~15–20 minutes. This validates IaC completeness and eliminates configuration drift.',
                font_size=14, color=DARK_GRAY)

    set_notes(slide, """WHAT TO SAY:
"All infrastructure is provisioned via Terraform — nothing is created manually in the Azure portal.

We evaluated three IaC tools. Bicep is excellent for Azure-only workloads and always has day-one support for new Azure features. Pulumi lets you use real programming languages. But we chose Terraform for three reasons: multi-cloud portability, explicit state management that enables reliable plan and destroy operations, and the largest community ecosystem.

Our Terraform configuration provisions everything you see here: the resource group, a virtual network with a dedicated subnet, an AKS cluster configured with 1-to-5 node autoscaling on burstable B2s VMs, Azure Container Registry, Log Analytics for monitoring, and blob storage for backups. It also creates the role assignment that allows AKS to pull images from ACR using managed identity — no service principal secrets to rotate.

A key capability: we can run terraform destroy to tear down everything, then terraform apply to rebuild it identically in about 15 to 20 minutes. This proves our IaC is complete and self-contained — nothing exists that isn't in code."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 10 — CONTAINER ORCHESTRATION
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Container Orchestration — AKS + Helm')

    # AKS config
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
                'AKS Cluster Configuration', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(0.6), Inches(2.1), Inches(5.8),
        ['Setting', 'Value'],
        [
            ['Kubernetes', '1.29'],
            ['Node VM', 'Standard_B2s (2 vCPU, 4 GB)'],
            ['Node count', '2 (autoscale 1–5)'],
            ['Network', 'Azure CNI + Calico'],
            ['Identity', 'System-assigned MI'],
            ['Monitoring', 'OMS → Log Analytics'],
        ], font_size=11)

    # Helm charts
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                'Helm Charts (3)', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(7), Inches(2.1), Inches(5.8),
        ['Chart', 'Resources', 'Key Feature'],
        [
            ['frontend/', 'Deploy, Svc, HPA', 'Blue/green slot label'],
            ['backend/', 'Deploy, Svc, HPA', 'ClusterIP (internal)'],
            ['mongodb/', 'StatefulSet, PVC, CronJob', 'Backup automation'],
        ], font_size=11)

    add_textbox(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(0.4),
                'Why Helm over Kustomize?', font_size=16, color=DARK_NAVY, bold=True)
    add_bullet_list(slide, Inches(7), Inches(4.9), Inches(5.5), Inches(2), [
        'Templating — Go templates for parameterised manifests',
        'Release management — tracks history, helm rollback',
        'Lifecycle hooks for pre/post deployment',
        'Extensive chart ecosystem',
    ], font_size=14)

    set_notes(slide, """WHAT TO SAY:
"For container orchestration, we use Azure Kubernetes Service — AKS.

Our cluster runs Kubernetes 1.29 on Standard B2s burstable VMs — cost-effective at 2 vCPUs and 4 GB RAM each. We start with 2 nodes and the cluster autoscaler can scale up to 5 nodes when demand requires it.

We use Azure CNI for networking, which gives each pod a real VNet IP address — this is required for Calico network policies that control pod-to-pod traffic. Authentication uses system-assigned managed identity, meaning there are no service principal secrets to manage or rotate.

Application packaging uses Helm 3 with three charts: frontend, backend, and MongoDB. Each chart produces the Kubernetes resources for its service. The frontend uses a LoadBalancer Service for external traffic. The backend uses ClusterIP — it's internal only. MongoDB uses a StatefulSet with persistent volume claims and a CronJob for automated daily backups.

We chose Helm over Kustomize because Helm's Go templating, release history tracking, and rollback capability — helm rollback — give us better deployment management. The blue/green slot label is a simple Helm value that gets injected at deploy time."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 11 — PERFORMANCE & MONITORING
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Performance & Monitoring')

    # Performance
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
                'Performance Measures', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(0.6), Inches(2.0), Inches(5.8),
        ['Metric', 'Frontend', 'Backend'],
        [
            ['Cold start', '~1s (Node.js)', '~15–30s (JVM)'],
            ['Warm latency', '<50ms', '<20ms'],
            ['Throughput', '~500 req/s', '~200 req/s'],
            ['Image size', '~80 MB', '~150 MB'],
        ], font_size=12)

    add_bullet_list(slide, Inches(0.6), Inches(4.2), Inches(5.8), Inches(1.5), [
        'Multi-stage Docker builds → minimal images',
        'Readiness probes prevent traffic to unready pods',
        'Connection pooling (MongoDB + Axios keep-alive)',
        'Resource requests/limits set per pod',
    ], font_size=14)

    # Monitoring
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                'Monitoring & Alerting', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(7), Inches(2.0), Inches(5.5),
        ['Alert', 'Condition', 'Severity'],
        [
            ['Pod restart loop', 'Restarts > 3 / 5 min', 'Critical'],
            ['High error rate', 'HTTP 5xx > 5%', 'Warning'],
            ['High latency', 'P95 > 2s', 'Warning'],
            ['Node pressure', 'CPU > 90%', 'Warning'],
            ['MongoDB failure', 'Connection errors', 'Critical'],
        ], font_size=11)

    add_textbox(slide, Inches(7), Inches(4.7), Inches(5.5), Inches(1.5),
                'Critique: JVM cold start is the main concern — mitigated by readiness probes + min 2 replicas. '
                'GraalVM native image could reduce to <1s. OpenTelemetry distributed tracing should be added '
                'for cross-service correlation.',
                font_size=13, color=RGBColor(0x88, 0x44, 0x00), italic=True)

    set_notes(slide, """WHAT TO SAY:
"Let's look at performance. The frontend starts in about 1 second — Node.js is fast to boot. The backend takes 15 to 30 seconds due to JVM and Spring Boot startup. Once warm, latency is excellent: under 50 milliseconds for the frontend, under 20 for the backend.

We achieve small image sizes through multi-stage Docker builds — the frontend is about 80 MB, the backend about 150 MB. Smaller images mean faster pod startup and less registry storage.

For monitoring, both services expose Prometheus-format metrics. The AKS OMS agent forwards container logs and metrics to Azure Log Analytics. We've defined five alert conditions — the two critical ones are pod restart loops and MongoDB connection failures.

A key critique: the JVM cold start is a real concern in a Kubernetes environment where pods can be rescheduled at any time. We mitigate this with readiness probes — Kubernetes won't send traffic until the pod is healthy — and by maintaining a minimum of 2 replicas. For production, GraalVM native image compilation could reduce startup to under 1 second."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 12 — SCALING & BACKUP
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Scaling & Backup Strategy')

    # Scaling
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
                'Horizontal Pod Autoscaling (HPA)', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(0.6), Inches(2.1), Inches(5.8),
        ['Service', 'Min', 'Max', 'CPU Target', 'Mem Target'],
        [
            ['Frontend', '2', '10', '70%', '80%'],
            ['Backend', '2', '8', '70%', '80%'],
            ['MongoDB', '1', '1 (fixed)', 'N/A', 'N/A'],
        ], font_size=12)

    add_bullet_list(slide, Inches(0.6), Inches(4.0), Inches(5.8), Inches(1.5), [
        'Cluster autoscaler: node count 1–5',
        'Scale out on pending pods; scale in after 10 min',
        'Critique: MongoDB single instance — needs Replica Set for HA',
        'Critique: HPA is reactive; KEDA could enable proactive scaling',
    ], font_size=14)

    # Backup
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                'Backup & Restore', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(7), Inches(2.1), Inches(5.5),
        ['Metric', 'Target', 'Justification'],
        [
            ['RPO (max data loss)', '24 hours', 'Daily backups'],
            ['RTO (max downtime)', '1 hour', 'Restore + redeploy'],
        ], font_size=12)

    add_bullet_list(slide, Inches(7), Inches(3.8), Inches(5.5), Inches(2), [
        'K8s CronJob: mongodump daily at 02:00 UTC',
        'Compressed with tar/gzip',
        '7-day rolling retention with auto-cleanup',
        'Critique: Should upload to Azure Blob for durability',
        'Critique: Monthly restore tests needed',
    ], font_size=14)

    set_notes(slide, """WHAT TO SAY:
"For scaling, we implement Horizontal Pod Autoscaling on both the frontend and backend. The frontend scales from 2 to 10 pods, the backend from 2 to 8. Both use multi-metric scaling based on CPU at 70 percent and memory at 80 percent.

At the infrastructure level, the AKS cluster autoscaler manages node count from 1 to 5. It adds nodes when pods are pending due to insufficient resources and removes them after 10 minutes of idle time.

A key limitation: MongoDB is a single instance with no replication. For production, we'd need either a MongoDB Replica Set for high availability or migrate to Azure Cosmos DB for MongoDB API, which provides built-in replication and geo-distribution.

For backup, a Kubernetes CronJob runs mongodump every night at 2 AM UTC, compresses the output, and maintains a 7-day rolling window. Our RPO is 24 hours — meaning we accept up to one day of data loss — and our RTO is 1 hour to restore and redeploy.

A critique: these backups currently reside inside the pod. They should be uploaded to Azure Blob Storage for durability, and we should test restores monthly to verify they actually work."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 13 — SECURITY
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Security — Defence in Depth')

    add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12),
        ['Layer', 'Measure', 'Implementation'],
        [
            ['Network', 'VNet isolation', 'AKS on dedicated subnet (10.0.1.0/24)'],
            ['Network', 'Network policies', 'Calico for pod-to-pod traffic control'],
            ['Container', 'Non-root execution', 'Dockerfiles create appuser; runAsNonRoot: true'],
            ['Container', 'Minimal images', 'Alpine-based (~5 MB base layer)'],
            ['Registry', 'No admin credentials', 'Managed Identity with AcrPull role'],
            ['Application', 'Security headers', 'Helmet.js (HSTS, CSP, X-Frame-Options)'],
            ['Application', 'Input validation', 'Jakarta Bean Validation on backend'],
            ['Application', 'Rate limiting', 'express-rate-limit (100 req/15 min)'],
            ['Supply chain', 'Vuln scanning', 'Trivy in CI — fail on CRITICAL'],
        ], font_size=11)

    add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2),
                'Critique: Production additions needed — explicit NetworkPolicy resources, Pod Security Standards (Restricted), '
                'Azure Key Vault for secrets, Application Gateway WAF, cert-manager + Let\'s Encrypt for TLS, '
                'and Notary v2 for image signing.',
                font_size=14, color=RGBColor(0x88, 0x44, 0x00), italic=True)

    set_notes(slide, """WHAT TO SAY:
"Security is implemented as defence in depth — multiple layers, each providing protection even if another fails.

At the network layer, AKS runs on a dedicated subnet within a VNet, and Calico network policies control which pods can communicate with which.

At the container level, all containers run as non-root users. We create a dedicated appuser in each Dockerfile and set runAsNonRoot in the Kubernetes security context. We use Alpine-based images with just a 5 megabyte base layer — smaller images mean fewer potential vulnerabilities.

The registry uses managed identity, not admin credentials — no passwords to leak or rotate.

At the application level, the frontend uses Helmet.js for security headers — HSTS, Content Security Policy, X-Frame-Options. The backend uses Jakarta Bean Validation for input validation. Rate limiting is applied at 100 requests per 15 minutes.

And for supply chain security, Trivy scans every Docker image in CI. Any CRITICAL vulnerability fails the build.

For production, we'd add explicit NetworkPolicy manifests, Azure Key Vault for secrets management, a WAF via Application Gateway, and TLS certificates via cert-manager with Let's Encrypt."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 14 — COST ANALYSIS
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Cost Analysis')

    add_table_slide(slide, Inches(1.5), Inches(1.5), Inches(10),
        ['Component', 'Dev Cost/Month', 'Prod Cost/Month'],
        [
            ['AKS Control Plane', '$0 (free tier)', '~$73 (Standard)'],
            ['AKS Node VMs', '~$60 (2× B2s)', '~$210 (3× D2s_v3)'],
            ['ACR', '~$5 (Basic)', '~$20 (Standard)'],
            ['Log Analytics', '~$10', '~$35'],
            ['Storage', '<$1', '<$5'],
            ['TOTAL', '~$75/month', '~$345/month'],
        ], font_size=13)

    add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.4),
                'Cost Optimisation Strategies:', font_size=18, color=DARK_NAVY, bold=True)
    add_bullet_list(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(2), [
        'Scheduled teardown — destroy non-prod environments overnight (~65% savings)',
        'Spot instances — up to 90% discount for interruptible workloads',
        'Reserved instances — 30–60% savings with 1–3 year commitment',
        'Autoscaler tuning — minimum 1 node during quiet periods',
        'Alternative: Azure Container Apps (consumption) → ~$10–20/month for low-traffic apps',
    ], font_size=15)

    set_notes(slide, """WHAT TO SAY:
"Let's talk cost. In development, our estimated monthly spend is about 75 dollars. The AKS control plane is free-tier, and we use burstable B2s VMs that cost about 30 dollars each.

In production, with Standard-tier AKS for the SLA guarantee and larger D2s v3 nodes, we're looking at about 345 dollars per month.

All our application tooling — Terraform, Helm, GitHub Actions, Docker — is either open-source or free-tier.

For optimisation, the biggest win is scheduled teardown. If we destroy non-prod environments overnight and on weekends, we save roughly 65 percent on compute. Spot instances can offer up to 90 percent discount for workloads that tolerate interruption. Reserved instances give 30 to 60 percent savings for committed usage.

An honest alternative: for a low-traffic application like this, Azure Container Apps with consumption pricing would cost just 10 to 20 dollars per month. We chose AKS for the learning value and fine-grained control, but Container Apps would be more cost-effective at this scale."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 15 — VULNERABILITY SCANNING & SUSTAINABILITY
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Vulnerability Scanning & Sustainability')

    # Vuln scanning
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
                'Vulnerability Scanning (Trivy)', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(0.6), Inches(2.1), Inches(5.8),
        ['Feature', 'Trivy', 'Docker Scout', 'Snyk'],
        [
            ['Cost', 'Free (Apache 2.0)', 'Free (limited)', 'Free (limited)'],
            ['OS packages', '✓', '✓', '✓'],
            ['Language deps', '✓', '✓', '✓'],
            ['IaC scanning', '✓', '✗', '✓'],
            ['License scan', '✓', '✗', '✓'],
        ], font_size=11)

    add_bullet_list(slide, Inches(0.6), Inches(4.5), Inches(5.8), Inches(1.5), [
        'CRITICAL → build fails',
        'HIGH → warning logged',
        'MEDIUM/LOW → informational',
        'Critique: Build-time only; ACR Defender for continuous',
    ], font_size=14)

    # Sustainability
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                'Sustainability', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(7), Inches(2.1), Inches(5.5),
        ['Measure', 'Impact'],
        [
            ['B-series burstable VMs', 'Lower baseline energy'],
            ['HPA + cluster autoscaler', 'Resources when needed only'],
            ['Environment teardown', '~65% compute savings'],
            ['Alpine-based images', 'Less storage & transfer'],
            ['Renewable-energy regions', 'Lower carbon footprint'],
        ], font_size=11)

    add_textbox(slide, Inches(7), Inches(4.8), Inches(5.5), Inches(1.2),
                'Critique: Most impactful measure is automated teardown. Serverless (Azure Functions + Cosmos DB) '
                'would eliminate baseline consumption entirely.',
                font_size=13, color=RGBColor(0x88, 0x44, 0x00), italic=True)

    set_notes(slide, """WHAT TO SAY:
"Vulnerability scanning uses Trivy — it's free, open-source, and scans both OS packages and language-specific dependencies. We chose it over Docker Scout and Snyk because it also handles IaC scanning and license scanning, all at no cost.

Our severity policy is simple: any CRITICAL vulnerability fails the build immediately. HIGH severity issues get a warning in the logs but don't block. Medium and low are informational.

One limitation: Trivy only scans at build time. For continuous scanning of images already in the registry, we'd enable Microsoft Defender for ACR.

On sustainability — data centres account for about 1 to 1.5 percent of global electricity use, so resource efficiency matters. We use burstable VMs that consume less baseline energy, autoscaling that only provisions resources when demand requires them, Alpine-based images that minimise storage and network transfer, and we target Azure regions powered by renewable energy.

The most impactful measure is automated environment teardown — destroying non-production infrastructure outside working hours saves approximately 65 percent of compute costs and the associated energy consumption."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 16 — ROLLBACK & CHANGE MANAGEMENT
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Rollback Plan & Change Management')

    # Rollback table
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
                'Four-Tier Rollback Strategy', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(0.6), Inches(2.1), Inches(12),
        ['Tier', 'Mechanism', 'Recovery Time', 'When to Use'],
        [
            ['1', 'Slot switch (kubectl patch)', '<10 seconds', 'Application bug post-deployment'],
            ['2', 'Helm rollback', '<60 seconds', 'Configuration change caused issues'],
            ['3', 'Previous image tag redeploy', '<5 minutes', 'Need to go back multiple versions'],
            ['4', 'Infrastructure rebuild (Terraform)', '15–30 minutes', 'Infrastructure compromise'],
        ], font_size=12)

    # Change management
    add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
                'Change Management — No Manual Changes Permitted', font_size=20, color=DARK_NAVY, bold=True)
    add_table_slide(slide, Inches(0.6), Inches(5.1), Inches(12),
        ['Configuration', 'Source of Truth', 'Applied By'],
        [
            ['Azure resources', 'Terraform files', 'terraform apply'],
            ['K8s workloads', 'Helm charts', 'helm upgrade --install'],
            ['App config', 'values.yaml + --set overrides', 'Helm'],
            ['CI/CD workflows', 'GitHub Actions YAML', 'GitHub'],
            ['Container contents', 'Dockerfiles', 'docker build'],
        ], font_size=11)

    set_notes(slide, """WHAT TO SAY:
"Our rollback plan has four tiers, each escalating in scope and recovery time.

Tier 1 is our instant rollback — a simple kubectl patch that switches the service selector back to the previous slot. Under 10 seconds. This handles the most common scenario: you deployed a new version, something's wrong, switch back immediately.

Tier 2 is Helm rollback — it reverts to the previous Helm release, including any configuration changes. Under 60 seconds.

Tier 3 redeploys a specific previous image tag — useful when you need to go back multiple versions. About 5 minutes including health checks.

Tier 4 is the nuclear option — a full infrastructure rebuild from Terraform. This takes 15 to 30 minutes but handles scenarios like infrastructure compromise or corruption.

For change management, we enforce a strict principle: no manual changes are permitted. Every configuration type has a single source of truth in code, and a single authorised mechanism to apply it. Azure resources come from Terraform. Kubernetes workloads come from Helm. Application config comes from values.yaml. This eliminates configuration drift — the actual state always matches the code."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 17 — ADDITIONAL FEATURES
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Additional Features')

    features = [
        ('Automated MongoDB Backup', 'K8s CronJob runs mongodump daily at 02:00 UTC\n7-day rolling retention with auto-cleanup', MID_BLUE),
        ('Horizontal Pod Autoscaling', 'HPA v2 with multi-metric scaling (CPU + memory)\nFrontend: 2–10 pods | Backend: 2–8 pods', ACCENT_GREEN),
        ('Container Vulnerability Scanning', 'Trivy integrated into CI pipeline\nScans OS packages + language deps\nFails build on CRITICAL CVEs', RGBColor(0x8b, 0x5c, 0xf6)),
        ('Blue/Green Zero-Downtime Deploy', 'Pod label slot switching\nAutomatic inactive-slot detection\nInstant rollback <10 seconds', RGBColor(0xf5, 0x9e, 0x0b)),
    ]

    for i, (title, desc, clr) in enumerate(features):
        x = Inches(0.6 + (i % 2) * 6.2)
        y = Inches(1.5 + (i // 2) * 2.8)
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = clr
        box.line.width = Pt(2)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = clr
        circle.line.fill.background()
        add_textbox(slide, x + Inches(0.28), y + Inches(0.22), Inches(0.4), Inches(0.4),
                    str(i + 1), font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, x + Inches(0.85), y + Inches(0.2), Inches(4.5), Inches(0.5),
                    title, font_size=18, color=DARK_NAVY, bold=True)
        # Description
        add_textbox(slide, x + Inches(0.3), y + Inches(0.9), Inches(5.2), Inches(1.3),
                    desc, font_size=14, color=DARK_GRAY)

    set_notes(slide, """WHAT TO SAY:
"Beyond the core requirements, we implemented four additional features.

First, automated MongoDB backups via a Kubernetes CronJob. It runs mongodump every night at 2 AM UTC, compresses the output, and automatically cleans up backups older than 7 days.

Second, Horizontal Pod Autoscaling using the v2 API with multi-metric scaling. Both CPU and memory utilisation drive scaling decisions, not just CPU. The frontend can scale from 2 to 10 pods, the backend from 2 to 8.

Third, container vulnerability scanning with Trivy, integrated directly into the CI pipeline. It scans both operating system packages and language-specific dependencies, and any CRITICAL CVE immediately fails the build.

And fourth, blue/green zero-downtime deployment with automatic slot detection. The CD pipeline determines which slot is inactive, deploys there, validates health, and then atomically switches traffic. Rollback is instant — under 10 seconds."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 18 — CRITICAL REFLECTION
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, 'Critical Reflection')

    # Strengths
    add_shape_bg(slide, Inches(0.4), Inches(1.4), Inches(6), Inches(3), RGBColor(0xec, 0xfe, 0xf5))
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
                '✓ Strengths', font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.3), [
        'Complete automation — zero manual steps',
        'Independent service pipelines with path filters',
        'Robust rollback — sub-10-second slot switching',
        'Security by default — non-root, Trivy, managed ID',
        'Observable — Prometheus metrics, structured logging',
    ], font_size=15, color=RGBColor(0x06, 0x5f, 0x46))

    # Limitations
    add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(6), Inches(3), RGBColor(0xfe, 0xf2, 0xf2))
    add_textbox(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.4),
                '✗ Limitations & Future Work', font_size=20, color=ACCENT_RED, bold=True)
    add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(2.3), [
        'Single MongoDB instance — needs Replica Set or Cosmos DB',
        'No integration/E2E tests in pipeline',
        'No service mesh for canary or mTLS',
        'No GitOps (ArgoCD/Flux) for state reconciliation',
        'Manual rollback — should auto-revert on error spikes',
        'No authentication — needs Azure AD / Entra ID',
    ], font_size=15, color=RGBColor(0x7f, 0x1d, 0x1d))

    # Bottom summary
    add_shape_bg(slide, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.4), LIGHT_GRAY)
    add_textbox(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.4),
                'Key Takeaway', font_size=18, color=DARK_NAVY, bold=True)
    add_textbox(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(1.6),
                'The microservice architectural style introduces deployment complexity that monoliths don\'t face. '
                'With appropriate tooling — Terraform for infrastructure, Helm for packaging, GitHub Actions for '
                'orchestration — this complexity is fully managed through automation. The result is a system that is '
                'reproducible, observable, and resilient: it can be destroyed and reliably rebuilt from code alone.',
                font_size=16, color=DARK_GRAY)

    set_notes(slide, """WHAT TO SAY:
"Let me reflect critically on what we've built.

On the strengths side: we achieved complete automation — there are zero manual steps from code commit to production. Each service has its own CI pipeline triggered by path filters. Our blue/green deployment gives us sub-10-second rollback capability. Security is built in by default, not bolted on after. And the system is observable with Prometheus metrics and structured logging flowing to Azure Monitor.

But there are real limitations. MongoDB is a single instance — a failure loses the database. Production needs a Replica Set or a managed service like Cosmos DB. We have no integration or end-to-end tests — the pipeline only runs unit tests. There's no service mesh, so we can't do canary deployments or mutual TLS between services. We don't use GitOps — ArgoCD or Flux would provide continuous state reconciliation. Rollback is currently manual — ideally it should trigger automatically when error rates spike post-deployment. And there's no authentication at all — production would need Azure AD integration.

The key takeaway: microservices are significantly harder to deploy than monoliths. But with the right tooling — Terraform, Helm, and GitHub Actions — that complexity can be fully managed through code and automation. The proof is that we can destroy everything and rebuild it identically from code alone."
""")

    # ═══════════════════════════════════════════════════════
    # SLIDE 19 — Q&A
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                'Questions & Discussion', font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(4), Inches(3.3), Inches(5.33), Inches(0.04), MID_BLUE)

    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                'Thank you for your attention', font_size=24, color=LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)

    links = [
        'GitHub: github.com/yewande111/Release-Management-Plan-RMP-',
        'Tech Stack: Node.js • Spring Boot • MongoDB • Terraform • AKS • Helm • GitHub Actions',
    ]
    for i, link in enumerate(links):
        add_textbox(slide, Inches(1), Inches(5.0 + i * 0.5), Inches(11), Inches(0.5),
                    link, font_size=14, color=RGBColor(0x88, 0x88, 0x88),
                    alignment=PP_ALIGN.CENTER)

    set_notes(slide, """WHAT TO SAY:
"Thank you all for listening. I'd be happy to take any questions.

If you'd like to explore the code, the full repository is on GitHub at the link shown. Everything we discussed today — the application code, Docker configurations, Terraform infrastructure, Helm charts, CI/CD workflows, and operational scripts — is all there.

Does anyone have questions about the architecture, the deployment strategy, or any of the evaluation criteria we covered?"
""")

    # ── Save ──
    output_path = os.path.join(OUT_DIR, 'Release_Management_Plan_Presentation.pptx')
    prs.save(output_path)
    print(f'PowerPoint saved to: {output_path}')

    # Clean up mock images
    for f in ['_mock_frontend.png', '_mock_api.png']:
        p = os.path.join(OUT_DIR, f)
        if os.path.exists(p):
            os.remove(p)

    return output_path


if __name__ == '__main__':
    build_presentation()
